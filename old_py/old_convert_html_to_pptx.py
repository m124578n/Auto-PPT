# html_to_pptx_accurate.py
# 精准的 HTML → PPTX 转换器

import os
import re
from typing import Dict, List, Tuple

from bs4 import BeautifulSoup
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class AccurateHTMLToPPTXConverter:
    """精准的 HTML 转 PPTX 转换器"""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)

    def _extract_image_path(self, img_src: str) -> str:
        """从 img src 提取实际的图片路径"""
        # src 格式: "downloaded_images/xxxxx.jpg"
        if img_src and os.path.exists(img_src):
            return img_src
        return None

    def _calculate_image_size(self, image_path: str, max_width: float, max_height: float) -> tuple:
        """计算保持宽高比的图片尺寸
        
        Args:
            image_path: 图片路径
            max_width: 最大宽度（英寸）
            max_height: 最大高度（英寸）
        
        Returns:
            (width, height) 元组（英寸）
        """
        try:
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
                
                # 计算基于宽度和高度的两种可能尺寸
                width_based = max_width
                height_based = max_width / aspect_ratio
                
                height_limit = max_height
                width_limit = max_height * aspect_ratio
                
                # 选择不超出边界的最大尺寸
                if height_based <= max_height:
                    return (width_based, height_based)
                else:
                    return (width_limit, height_limit)
        except Exception as e:
            print(f"   ⚠️ 无法读取图片尺寸：{e}")
            # 返回默认尺寸
            return (max_width, max_height)

    def parse_html(self, html_file: str) -> None:
        """解析 HTML 并转换为 PPTX"""
        print(f"\n📂 读取 HTML 文件：{html_file}")

        with open(html_file, "r", encoding="utf-8") as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, "html.parser")

        # 找到 presentation-container
        container = soup.find("div", class_="presentation-container")
        if not container:
            print("❌ 找不到 presentation-container")
            return

        # 找到所有 slide
        slides = container.find_all("div", class_="slide")
        print(f"   ✓ 找到 {len(slides)} 张幻灯片")

        for i, slide in enumerate(slides, 1):
            print(f"\n📝 处理第 {i} 张幻灯片...")

            try:
                # 判断 slide 类型
                slide_classes = slide.get("class", [])

                if "slide-opening" in slide_classes:
                    self._parse_opening(slide)
                    print(f"   ✓ 开场页创建成功")
                elif "slide-section" in slide_classes:
                    self._parse_section(slide)
                    print(f"   ✓ 章节页创建成功")
                elif "slide-closing" in slide_classes:
                    self._parse_closing(slide)
                    print(f"   ✓ 结尾页创建成功")
                else:
                    # 默认为内容页
                    self._parse_content(slide)
                    print(f"   ✓ 内容页创建成功")

            except Exception as e:
                print(f"   ❌ 创建失败：{e}")
                import traceback

                traceback.print_exc()

    def _parse_opening(self, slide_elem) -> None:
        """解析开场页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 背景渐层 - HTML: linear-gradient(135deg, #667eea 0%, #764ba2 100%)
        # 使用中间色 #6E72C6
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(110, 114, 198)
        bg.line.fill.background()

        # 提取主标题
        main_title = slide_elem.find("h1", class_="main-title")
        if main_title:
            title_text = main_title.get_text(strip=True)
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 根据标题长度调整字体大小
            if len(title_text) > 15:
                title_frame.paragraphs[0].font.size = Pt(52)
            else:
                title_frame.paragraphs[0].font.size = Pt(58)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            title_frame.word_wrap = True

        # 提取副标题
        subtitle = slide_elem.find("p", class_="subtitle")
        if subtitle:
            subtitle_text = subtitle.get_text(strip=True)
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.3), Inches(9), Inches(1.2)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 根据副标题长度调整字体大小
            if len(subtitle_text) > 25:
                subtitle_frame.paragraphs[0].font.size = Pt(26)
            else:
                subtitle_frame.paragraphs[0].font.size = Pt(28)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtitle_frame.word_wrap = True

    def _parse_section(self, slide_elem) -> None:
        """解析章节页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 背景 - HTML: linear-gradient(135deg, #4682b4 0%, #2c5f8d 100%)
        # 使用中间色 #3970A1
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(57, 112, 161)
        bg.line.fill.background()

        # 提取章节标题
        section_title = slide_elem.find("h2", class_="section-title")
        if section_title:
            title_text = section_title.get_text(strip=True)
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 根据标题长度调整字体大小
            if len(title_text) > 20:
                title_frame.paragraphs[0].font.size = Pt(46)
            else:
                title_frame.paragraphs[0].font.size = Pt(50)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            title_frame.word_wrap = True

        # 装饰线（上）
        line_top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.6), Inches(2), Inches(0.04)
        )
        line_top.fill.solid()
        line_top.fill.fore_color.rgb = RGBColor(255, 255, 255)
        line_top.line.fill.background()

        # 装饰线（下）
        line_bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.9), Inches(2), Inches(0.04)
        )
        line_bottom.fill.solid()
        line_bottom.fill.fore_color.rgb = RGBColor(255, 255, 255)
        line_bottom.line.fill.background()

    def _parse_closing(self, slide_elem) -> None:
        """解析结尾页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 背景 - HTML: linear-gradient(135deg, #f093fb 0%, #f5576c 100%)
        # 使用中间色 #F375B4
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(243, 117, 180)
        bg.line.fill.background()

        # 提取结尾标题（注意：HTML中是 h1）
        closing_title = slide_elem.find("h1", class_="closing-title")
        if closing_title:
            title_text = closing_title.get_text(strip=True)
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(9), Inches(1.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = title_text
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 根据标题长度调整字体大小
            if len(title_text) > 25:
                title_frame.paragraphs[0].font.size = Pt(46)
            else:
                title_frame.paragraphs[0].font.size = Pt(50)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            title_frame.word_wrap = True

        # 提取副文字
        subtext = slide_elem.find("p", class_="closing-subtext")
        if subtext:
            subtext_text = subtext.get_text(strip=True)
            subtext_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(1.0)
            )
            subtext_frame = subtext_box.text_frame
            subtext_frame.text = subtext_text
            subtext_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 根据副文字长度调整字体大小
            if len(subtext_text) > 20:
                subtext_frame.paragraphs[0].font.size = Pt(24)
            else:
                subtext_frame.paragraphs[0].font.size = Pt(26)
            subtext_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            subtext_frame.word_wrap = True

    def _parse_content(self, slide_elem) -> None:
        """解析内容页（包括纯文本和图文混合）"""
        # 检查是否有图片容器
        has_image_text = slide_elem.find("div", class_="image-text-container")
        has_full_image = slide_elem.find("div", class_="full-image-container")

        if has_full_image:
            self._parse_full_image(slide_elem)
        elif has_image_text:
            self._parse_image_text(slide_elem)
        else:
            self._parse_text_content(slide_elem)

    def _parse_text_content(self, slide_elem) -> None:
        """解析纯文本内容页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 提取标题
        title_elem = slide_elem.find("h2", class_="slide-title")
        if title_elem:
            title = title_elem.get_text(strip=True)
            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.6), Inches(8.8), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(38)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50
            title_frame.word_wrap = True

            # 标题下划线 - HTML: 4px solid #4682b4
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(8.8), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(70, 130, 180)  # #4682b4
            line.line.fill.background()

        # 提取项目符号列表
        bullet_list = slide_elem.find("ul", class_="bullet-list")
        if bullet_list:
            items = bullet_list.find_all("li")

            if items:
                # 扩大文字框高度，确保不会超出
                content_box = slide.shapes.add_textbox(
                    Inches(1.0), Inches(1.8), Inches(8.0), Inches(5.4)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                # 根据项目数量和内容长度调整字体大小和间距
                num_items = len(items)
                # 计算平均文字长度
                avg_length = sum(len(item.get_text(strip=True)) for item in items) / num_items
                
                if num_items >= 5 and avg_length > 25:
                    # 项目多且文字长，使用最小字体和间距
                    base_font_size = 21
                    indent_font_size = 19
                    base_spacing = 12
                    indent_spacing = 10
                    line_spacing = 1.25
                elif num_items >= 5:
                    # 项目较多，减小字体和间距
                    base_font_size = 22
                    indent_font_size = 20
                    base_spacing = 14
                    indent_spacing = 12
                    line_spacing = 1.3
                elif num_items >= 4:
                    # 中等数量
                    base_font_size = 23
                    indent_font_size = 21
                    base_spacing = 16
                    indent_spacing = 14
                    line_spacing = 1.35
                else:
                    # 项目较少，使用正常字体和间距
                    base_font_size = 24
                    indent_font_size = 22
                    base_spacing = 18
                    indent_spacing = 16
                    line_spacing = 1.4

                for i, item in enumerate(items):
                    # 判断是否为缩进项
                    is_indent = "indent-1" in item.get("class", [])
                    text = item.get_text(strip=True)

                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    if is_indent:
                        p.text = f"▸ {text}"
                        p.font.size = Pt(indent_font_size)
                        p.level = 1
                        p.space_after = Pt(indent_spacing)
                        p.font.color.rgb = RGBColor(85, 85, 85)  # #555
                    else:
                        p.text = f"▶ {text}"
                        p.font.size = Pt(base_font_size)
                        p.space_after = Pt(base_spacing)
                        p.font.color.rgb = RGBColor(52, 73, 94)  # #34495e
                    
                    p.line_spacing = line_spacing

    def _parse_image_text(self, slide_elem) -> None:
        """解析图文混合页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 提取标题
        title_elem = slide_elem.find("h2", class_="slide-title")
        if title_elem:
            title = title_elem.get_text(strip=True)
            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.75)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            # 根据标题长度调整字体大小
            if len(title) > 20:
                title_frame.paragraphs[0].font.size = Pt(34)
            else:
                title_frame.paragraphs[0].font.size = Pt(36)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50
            title_frame.word_wrap = True

        # 检查布局类型
        container = slide_elem.find("div", class_="image-text-container")
        is_horizontal = container and "layout-horizontal" in container.get("class", [])

        # 提取图片
        image_box = slide_elem.find("div", class_="image-box")
        if image_box:
            img_tag = image_box.find("img")
            if img_tag:
                img_src = img_tag.get("src", "")
                image_path = self._extract_image_path(img_src)

                if image_path and os.path.exists(image_path):
                    try:
                        if is_horizontal:
                            # 左图右文 - 保持宽高比
                            max_width = 4.4
                            max_height = 5.0
                            img_width, img_height = self._calculate_image_size(
                                image_path, max_width, max_height
                            )
                            # 居中对齐图片
                            left = Inches(0.7 + (max_width - img_width) / 2)
                            top = Inches(1.9 + (max_height - img_height) / 2)
                            
                            pic = slide.shapes.add_picture(
                                image_path,
                                left,
                                top,
                                width=Inches(img_width),
                                height=Inches(img_height),
                            )
                        else:
                            # 上图下文（vertical）- 保持宽高比
                            max_width = 7.0
                            max_height = 3.2
                            img_width, img_height = self._calculate_image_size(
                                image_path, max_width, max_height
                            )
                            # 居中对齐图片
                            left = Inches(1.5 + (max_width - img_width) / 2)
                            top = Inches(1.8 + (max_height - img_height) / 2)
                            
                            pic = slide.shapes.add_picture(
                                image_path,
                                left,
                                top,
                                width=Inches(img_width),
                                height=Inches(img_height),
                            )

                        print(f"   ✓ 添加图片：{os.path.basename(image_path)} ({img_width:.2f}\" × {img_height:.2f}\")")
                    except Exception as e:
                        print(f"   ⚠️ 图片添加失败：{e}")
                else:
                    print(f"   ⚠️ 图片不存在：{img_src}")

        # 提取文字
        text_box_elem = slide_elem.find("div", class_="text-box")
        if text_box_elem:
            text_content = text_box_elem.get_text(strip=True)

            if is_horizontal:
                # 右侧文字 - 扩大文字框以容纳更多内容
                text_box = slide.shapes.add_textbox(
                    Inches(5.35), Inches(1.8), Inches(4.0), Inches(5.4)
                )
                # 根据文字长度调整字体大小
                if len(text_content) > 200:
                    font_size = 19
                    line_spacing = 1.5
                elif len(text_content) > 150:
                    font_size = 20
                    line_spacing = 1.55
                else:
                    font_size = 21
                    line_spacing = 1.6
            else:
                # 下方文字 - 扩大文字框高度
                text_box = slide.shapes.add_textbox(
                    Inches(1), Inches(5.3), Inches(8), Inches(1.9)
                )
                # 根据文字长度调整字体大小
                if len(text_content) > 150:
                    font_size = 18
                    line_spacing = 1.4
                else:
                    font_size = 19
                    line_spacing = 1.5

            text_frame = text_box.text_frame
            text_frame.text = text_content
            text_frame.word_wrap = True
            text_frame.paragraphs[0].font.size = Pt(font_size)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50
            text_frame.paragraphs[0].line_spacing = line_spacing

    def _parse_full_image(self, slide_elem) -> None:
        """解析大图展示页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # 提取标题
        title_elem = slide_elem.find("h2", class_="slide-title")
        if title_elem:
            title = title_elem.get_text(strip=True)
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(8), Inches(0.75)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # 根据标题长度调整字体大小
            if len(title) > 20:
                title_frame.paragraphs[0].font.size = Pt(34)
            else:
                title_frame.paragraphs[0].font.size = Pt(36)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50
            title_frame.word_wrap = True

        # 提取图片
        full_image_container = slide_elem.find("div", class_="full-image-container")
        if full_image_container:
            img_tag = full_image_container.find("img")
            if img_tag:
                img_src = img_tag.get("src", "")
                image_path = self._extract_image_path(img_src)

                if image_path and os.path.exists(image_path):
                    try:
                        # 大图展示 - HTML: max-width 90%, max-height 70%
                        # 保持宽高比
                        max_width = 7.5
                        max_height = 4.2
                        img_width, img_height = self._calculate_image_size(
                            image_path, max_width, max_height
                        )
                        # 居中对齐图片
                        left = Inches(1.25 + (max_width - img_width) / 2)
                        top = Inches(1.9 + (max_height - img_height) / 2)
                        
                        pic = slide.shapes.add_picture(
                            image_path,
                            left,
                            top,
                            width=Inches(img_width),
                            height=Inches(img_height),
                        )
                        print(f"   ✓ 添加大图：{os.path.basename(image_path)} ({img_width:.2f}\" × {img_height:.2f}\")")
                    except Exception as e:
                        print(f"   ⚠️ 图片添加失败：{e}")
                else:
                    print(f"   ⚠️ 图片不存在：{img_src}")

            # 提取图片说明
            caption = full_image_container.find("p", class_="caption")
            if caption:
                caption_text = caption.get_text(strip=True)
                caption_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.7)
                )
                caption_frame = caption_box.text_frame
                caption_frame.text = caption_text
                caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                # 根据说明文字长度调整字体大小
                if len(caption_text) > 60:
                    caption_frame.paragraphs[0].font.size = Pt(15)
                else:
                    caption_frame.paragraphs[0].font.size = Pt(16)
                caption_frame.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)  # #7f8c8d
                caption_frame.word_wrap = True

    def save(self, output_path: str) -> None:
        """保存 PPTX"""
        self.prs.save(output_path)
        print(f"\n✅ PPTX 已保存：{output_path}")


# ==================== 主程序 ====================
def main():
    import sys

    print("🔄 精准 HTML → PPTX 转换器")
    print("=" * 60)

    # 寻找 HTML 文件
    html_files = [f for f in os.listdir(".") if f.endswith("_presentation.html")]

    if not html_files:
        print("❌ 找不到 HTML 演示文件")
        return

    # 如果有多个文件，列出选择
    if len(html_files) > 1:
        print("\n📋 找到多个 HTML 文件：")
        for i, f in enumerate(html_files, 1):
            print(f"   {i}. {f}")

        choice = input("\n请选择要转换的文件编号（直接按 Enter 选择第一个）: ").strip()

        if choice and choice.isdigit() and 1 <= int(choice) <= len(html_files):
            html_file = html_files[int(choice) - 1]
        else:
            html_file = html_files[0]
    else:
        html_file = html_files[0]

    print(f"\n📂 选择的文件：{html_file}")

    # 转换
    converter = AccurateHTMLToPPTXConverter()
    converter.parse_html(html_file)

    # 保存
    output_filename = html_file.replace("_presentation.html", ".pptx")
    converter.save(output_filename)

    # 验证
    if os.path.exists(output_filename):
        file_size = os.path.getsize(output_filename)
        print(f"\n📦 档案资讯：")
        print(f"   路径：{os.path.abspath(output_filename)}")
        print(f"   大小：{file_size:,} bytes ({file_size/1024/1024:.2f} MB)")

        try:
            verify_prs = Presentation(output_filename)
            print(f"   幻灯片数：{len(verify_prs.slides)}")
        except Exception as e:
            print(f"   ⚠️ 验证失败：{e}")

    print("\n" + "=" * 60)
    print("✅ 转换完成！")
    print(f"📁 请开启档案检查：{output_filename}")
    print("=" * 60)


if __name__ == "__main__":
    main()

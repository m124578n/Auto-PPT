"""
測試重構後的代碼

測試內容：
1. SlideType 註冊是否正常
2. HTML 生成是否正常
3. PPTX 生成是否正常
"""

import json
import os

from slide_generator import HTMLGenerator, PPTXGenerator
from slide_types import SlideTypeRegistry


def test_registry():
    """測試 Slide 類型註冊"""
    print("=" * 60)
    print("測試 1: Slide 類型註冊")
    print("=" * 60)
    
    types = SlideTypeRegistry.all_types()
    print(f"✓ 已註冊 {len(types)} 種 Slide 類型：")
    for t in types:
        print(f"  - {t}")
    
    # 測試獲取類型
    for t in types:
        slide_class = SlideTypeRegistry.get(t)
        assert slide_class is not None, f"無法獲取 {t} 類型"
    
    print("✓ 所有類型都可以正確獲取")
    print()


def test_html_generation():
    """測試 HTML 生成"""
    print("=" * 60)
    print("測試 2: HTML 生成")
    print("=" * 60)
    
    # 測試數據
    test_data = {
        'title': '測試簡報',
        'topic': 'test_presentation',
        'slides': [
            {
                'slide_type': 'opening',
                'title': '歡迎來到測試簡報',
                'subtitle': '這是重構後的版本'
            },
            {
                'slide_type': 'section_divider',
                'section_title': '第一章：架構設計'
            },
            {
                'slide_type': 'text_content',
                'title': '主要特點',
                'bullets': [
                    '使用 Strategy Pattern',
                    '易於擴展新的 Slide 類型',
                    'HTML 和 PPTX 生成邏輯統一'
                ],
                'indent_levels': [0, 0, 0]
            },
            {
                'slide_type': 'closing',
                'closing_text': '謝謝觀看',
                'subtext': '歡迎提供反饋'
            }
        ]
    }
    
    # 生成 HTML
    html_gen = HTMLGenerator()
    html_content = html_gen.generate_from_data(test_data)
    
    assert len(html_content) > 0, "HTML 內容為空"
    assert '<div class="slide slide-opening">' in html_content, "缺少開場頁"
    assert '<div class="slide slide-section">' in html_content, "缺少章節頁"
    assert '<div class="slide slide-content">' in html_content, "缺少內容頁"
    assert '<div class="slide slide-closing">' in html_content, "缺少結尾頁"
    
    print("✓ HTML 生成成功")
    
    # 保存測試 HTML
    output_file = 'test_refactored_presentation.html'
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"✓ HTML 已保存：{output_file}")
    print()


def test_pptx_generation():
    """測試 PPTX 生成"""
    print("=" * 60)
    print("測試 3: PPTX 生成")
    print("=" * 60)
    
    # 測試數據
    test_data = {
        'title': '測試簡報',
        'topic': 'test_presentation',
        'slides': [
            {
                'slide_type': 'opening',
                'title': '歡迎來到測試簡報',
                'subtitle': '這是重構後的版本'
            },
            {
                'slide_type': 'section_divider',
                'section_title': '第一章：架構設計'
            },
            {
                'slide_type': 'text_content',
                'title': '主要特點',
                'bullets': [
                    '使用 Strategy Pattern',
                    '易於擴展新的 Slide 類型',
                    'HTML 和 PPTX 生成邏輯統一',
                    'Registry Pattern 自動註冊'
                ],
                'indent_levels': [0, 0, 0, 1]
            },
            {
                'slide_type': 'text_content',
                'title': '如何擴展',
                'bullets': [
                    '繼承 SlideType 基類',
                    '實現 generate_html() 方法',
                    '實現 generate_pptx() 方法',
                    '使用 @SlideTypeRegistry.register() 裝飾器',
                    '完成！不需修改其他代碼'
                ],
                'indent_levels': [0, 1, 1, 0, 0]
            },
            {
                'slide_type': 'closing',
                'closing_text': '謝謝觀看',
                'subtext': '歡迎提供反饋'
            }
        ]
    }
    
    # 生成 PPTX
    pptx_gen = PPTXGenerator()
    prs = pptx_gen.generate_from_data(test_data)
    
    assert len(prs.slides) == len(test_data['slides']), f"Slide 數量不符：期望 {len(test_data['slides'])}，實際 {len(prs.slides)}"
    
    print(f"✓ PPTX 生成成功，共 {len(prs.slides)} 張幻燈片")
    
    # 保存測試 PPTX
    output_file = 'test_refactored.pptx'
    pptx_gen.save(output_file)
    
    # 驗證文件
    file_size = os.path.getsize(output_file)
    print(f"✓ PPTX 文件大小：{file_size:,} bytes ({file_size/1024:.2f} KB)")
    
    # 重新讀取驗證
    from pptx import Presentation
    verify_prs = Presentation(output_file)
    print(f"✓ 驗證成功：可以正常讀取 PPTX 文件")
    print()


def test_json_workflow():
    """測試 JSON → HTML/PPTX 工作流程"""
    print("=" * 60)
    print("測試 4: JSON 工作流程")
    print("=" * 60)
    
    # 創建測試 JSON
    test_data = {
        'title': '完整流程測試',
        'topic': 'workflow_test',
        'slides': [
            {
                'slide_type': 'opening',
                'title': '完整流程測試',
                'subtitle': 'JSON → HTML/PPTX'
            },
            {
                'slide_type': 'text_content',
                'title': '測試項目',
                'bullets': [
                    'JSON 數據保存',
                    'HTML 生成',
                    'PPTX 生成',
                    '文件驗證'
                ],
                'indent_levels': [0, 0, 0, 0]
            },
            {
                'slide_type': 'closing',
                'closing_text': '測試完成',
                'subtext': '所有功能正常'
            }
        ]
    }
    
    # 保存 JSON
    json_file = 'test_workflow_data.json'
    with open(json_file, 'w', encoding='utf-8') as f:
        json.dump(test_data, f, ensure_ascii=False, indent=2)
    
    print(f"✓ JSON 已保存：{json_file}")
    
    # 從 JSON 生成 HTML
    with open(json_file, 'r', encoding='utf-8') as f:
        loaded_data = json.load(f)
    
    html_gen = HTMLGenerator()
    html_content = html_gen.generate_from_data(loaded_data)
    
    html_file = 'test_workflow_presentation.html'
    with open(html_file, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"✓ HTML 已生成：{html_file}")
    
    # 從 JSON 生成 PPTX
    pptx_gen = PPTXGenerator()
    prs = pptx_gen.generate_from_data(loaded_data)
    
    pptx_file = 'test_workflow.pptx'
    pptx_gen.save(pptx_file)
    
    print(f"✓ PPTX 已生成：{pptx_file}")
    print()


def main():
    """運行所有測試"""
    print("\n" + "=" * 60)
    print("開始測試重構後的代碼")
    print("=" * 60)
    print()
    
    try:
        test_registry()
        test_html_generation()
        test_pptx_generation()
        test_json_workflow()
        
        print("=" * 60)
        print("✅ 所有測試通過！")
        print("=" * 60)
        print()
        print("生成的測試文件：")
        print("  - test_refactored_presentation.html")
        print("  - test_refactored.pptx")
        print("  - test_workflow_data.json")
        print("  - test_workflow_presentation.html")
        print("  - test_workflow.pptx")
        print()
        
    except AssertionError as e:
        print(f"\n❌ 測試失敗：{e}")
        import traceback
        traceback.print_exc()
    except Exception as e:
        print(f"\n❌ 發生錯誤：{e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()


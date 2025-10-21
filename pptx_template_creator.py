"""
PPTX 模板創建器
用 Gemini 分析現有 PPTX 並生成 JSON 模板文件
"""
import json
import os
from pathlib import Path
from typing import Dict, List, Optional

from google import genai
from google.genai import types
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from AutoPPT.utils.logger import get_logger

logger = get_logger()


class PPTXTemplateCreator:
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.client = genai.Client(api_key=api_key)
        
    def analyze_pptx(self, pptx_path: str) -> Dict:
        """分析 PPTX 文件結構"""
        logger.info(f"📄 開始分析 PPTX 文件：{pptx_path}")
        
        prs = Presentation(pptx_path)
        
        # 基本信息
        slide_width = prs.slide_width.inches
        slide_height = prs.slide_height.inches
        
        logger.info(f"   ✓ 尺寸：{slide_width}\" × {slide_height}\"")
        logger.info(f"   ✓ 總共 {len(prs.slides)} 張 slides")
        
        # 分析每張 slide
        slides_info = []
        for i, slide in enumerate(prs.slides, 1):
            slide_info = self._analyze_slide(slide, i)
            if slide_info:
                slides_info.append(slide_info)
        
        analysis = {
            "slide_width": slide_width,
            "slide_height": slide_height,
            "total_slides": len(prs.slides),
            "slides": slides_info
        }
        
        logger.info(f"   ✓ 分析完成：找到 {len(slides_info)} 種不同的佈局")
        
        return analysis
    
    def _analyze_slide(self, slide, slide_number: int) -> Dict:
        """分析單張 slide 的結構"""
        logger.info(f"\n📝 分析 Slide {slide_number}")
        
        # 提取所有元素
        elements = []
        
        for shape in slide.shapes:
            element = self._analyze_shape(shape)
            if element:
                elements.append(element)
                logger.info(f"   - {element['type']}: {element['name']}")
        
        # 提取背景顏色
        background = self._analyze_background(slide)
        
        return {
            "slide_number": slide_number,
            "layout_name": slide.slide_layout.name if hasattr(slide, 'slide_layout') else f"Layout_{slide_number}",
            "background": background,
            "elements": elements
        }
    
    def _analyze_shape(self, shape) -> Optional[Dict]:
        """分析單個 shape（文本框、圖片、形狀等）"""
        try:
            # 基本位置信息（轉換為英寸）
            position = {
                "left": round(shape.left.inches, 2) if hasattr(shape, 'left') else 0,
                "top": round(shape.top.inches, 2) if hasattr(shape, 'top') else 0,
                "width": round(shape.width.inches, 2) if hasattr(shape, 'width') else 0,
                "height": round(shape.height.inches, 2) if hasattr(shape, 'height') else 0
            }
            
            # 文本框
            if shape.has_text_frame:
                return self._analyze_textbox(shape, position)
            
            # 圖片
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return {
                    "type": "image",
                    "name": "image_placeholder",
                    "position": position
                }
            
            # 形狀（矩形、線條等）
            elif shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM]:
                return {
                    "type": "shape",
                    "name": f"shape_{shape.shape_type}",
                    "shape_type": "rectangle",  # 簡化處理
                    "position": position,
                    "style": self._analyze_shape_style(shape)
                }
            
            return None
            
        except Exception as e:
            logger.warning(f"⚠️  分析 shape 失敗：{e}")
            return None
    
    def _analyze_textbox(self, shape, position: Dict) -> Dict:
        """分析文本框"""
        text_frame = shape.text_frame
        
        # 獲取文本內容（用於識別用途）
        text_content = text_frame.text.strip() if text_frame.text else ""
        
        # 分析第一個段落的樣式
        style = {}
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            
            # 對齊方式
            if hasattr(p, 'alignment') and p.alignment is not None:
                alignment_map = {
                    1: "left",
                    2: "center",
                    3: "right"
                }
                style["alignment"] = alignment_map.get(int(p.alignment), "left")
            
            # 字體樣式
            if p.runs:
                run = p.runs[0]
                font = run.font
                
                if font.size:
                    style["font_size"] = int(font.size.pt)
                
                if font.bold:
                    style["font_bold"] = True
                
                if font.color and font.color.rgb:
                    rgb = font.color.rgb
                    style["font_color"] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}".upper()
        
        # 推測用途
        name = self._guess_textbox_purpose(text_content, position, style)
        
        return {
            "type": "textbox",
            "name": name,
            "position": position,
            "style": style,
            "sample_text": text_content[:50] if text_content else ""
        }
    
    def _guess_textbox_purpose(self, text: str, position: Dict, style: Dict) -> str:
        """根據位置和樣式推測文本框用途"""
        # 根據位置判斷
        top = position.get("top", 0)
        height = position.get("height", 0)
        font_size = style.get("font_size", 0)
        is_centered = style.get("alignment") == "center"
        is_bold = style.get("font_bold", False)
        
        # 標題（頂部、大字、粗體、居中）
        if top < 2 and font_size > 30 and is_bold and is_centered:
            return "title"
        
        # 副標題（靠近標題、中等大小、居中）
        elif top < 3 and font_size > 20 and is_centered:
            return "subtitle"
        
        # 內容（中間位置、較大區域）
        elif height > 3:
            return "content"
        
        # 說明文字（小字體）
        elif font_size < 20:
            return "caption"
        
        # 默認
        else:
            return "text"
    
    def _analyze_shape_style(self, shape) -> Dict:
        """分析形狀樣式"""
        style = {}
        
        try:
            if hasattr(shape, 'fill') and shape.fill:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                    rgb = shape.fill.fore_color.rgb
                    style["fill_color"] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}".upper()
        except:
            pass
        
        return style
    
    def _analyze_background(self, slide) -> Dict:
        """分析背景"""
        try:
            background = slide.background
            if hasattr(background, 'fill') and background.fill:
                if hasattr(background.fill, 'fore_color') and background.fill.fore_color.rgb:
                    rgb = background.fill.fore_color.rgb
                    return {
                        "type": "solid",
                        "color": f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}".upper()
                    }
        except:
            pass
        
        return {
            "type": "solid",
            "color": "#FFFFFF"
        }
    
    def generate_template_json(self, analysis: Dict, template_name: str = "自定義模板") -> str:
        """使用 Gemini 生成 JSON 模板"""
        logger.info("\n🤖 使用 Gemini 生成 JSON 模板...")
        
        # 準備 prompt
        prompt = self._create_gemini_prompt(analysis, template_name)
        
        # 調用 Gemini
        response = self.client.models.generate_content(
            model="gemini-2.0-flash-exp",
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.1,
                response_mime_type="application/json"
            )
        )
        
        logger.info("   ✓ Gemini 生成完成")
        
        return response.text
    
    def _create_gemini_prompt(self, analysis: Dict, template_name: str) -> str:
        """創建給 Gemini 的 prompt"""
        
        # 讀取現有模板作為參考
        reference_template = self._load_reference_template()
        
        prompt = f"""你是一個專業的 PowerPoint 模板設計師。我需要你根據分析的 PPTX 文件結構，生成一個符合我們系統格式的 JSON 模板。

# 目標
將以下 PPTX 分析結果轉換為 JSON 模板格式。

# PPTX 分析結果
```json
{json.dumps(analysis, ensure_ascii=False, indent=2)}
```

# 參考模板格式（請嚴格遵循這個結構）
```json
{reference_template}
```

# 生成要求

1. **template_info 部分**：
   - name: "{template_name}"
   - version: "1.0.0"
   - description: 根據分析的內容生成簡短描述
   - author: "Auto Generated"
   - slide_width 和 slide_height: 使用分析中的尺寸

2. **slide_types 部分**：
   - 為每個分析到的 slide 創建一個 slide type
   - type_id: 使用有意義的英文 ID（如 "opening", "content_1" 等）
   - name: 使用中文名稱（如 "開場頁", "內容頁 1" 等）
   - description: 簡短描述這個 slide 的用途
   - llm_instruction: 給 AI 的使用說明
   - json_schema: 定義這個 slide 需要的數據字段
   - pptx_layout: 包含 layout_index, background, elements

3. **元素處理**：
   - textbox: 保留 name, position, style
   - image: 如果是圖片佔位符，使用 max_width 和 max_height
   - shape: 保留裝飾性形狀

4. **樣式保留**：
   - 保留所有字體大小、顏色、對齊方式
   - 顏色使用十六進制格式（如 "#2C3E50"）

5. **智能命名**：
   - 根據文本框的用途給出有意義的 name
   - title, subtitle, content, caption 等

# 輸出要求
- 直接輸出 JSON，不要任何額外說明
- 確保 JSON 格式正確，可以直接解析
- 所有字符串使用中文（描述部分）
- 保持結構清晰、層次分明

請生成完整的 JSON 模板："""
        
        return prompt
    
    def _load_reference_template(self) -> str:
        """加載參考模板"""
        template_path = "templates/default_template.json"
        
        if os.path.exists(template_path):
            with open(template_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # 如果沒有參考模板，返回基本結構
        return """{
  "template_info": {
    "name": "示例模板",
    "version": "1.0.0",
    "description": "模板描述",
    "author": "作者",
    "slide_width": 10.0,
    "slide_height": 7.5
  },
  "slide_types": [
    {
      "type_id": "opening",
      "name": "開場頁",
      "description": "簡報的封面頁",
      "llm_instruction": "用於簡報開始",
      "json_schema": {
        "slide_type": "opening",
        "title": "標題",
        "subtitle": "副標題"
      },
      "pptx_layout": {
        "layout_index": 6,
        "background": {
          "type": "solid",
          "color": "#FFFFFF"
        },
        "elements": [
          {
            "type": "textbox",
            "name": "title",
            "position": {"left": 1.0, "top": 3.0, "width": 8.0, "height": 1.5},
            "style": {
              "font_size": 48,
              "font_bold": true,
              "font_color": "#000000",
              "alignment": "center"
            }
          }
        ]
      }
    }
  ]
}"""
    
    def create_template_from_pptx(
        self,
        pptx_path: str,
        output_path: Optional[str] = None,
        template_name: str = "自定義模板"
    ) -> str:
        """完整流程：從 PPTX 創建 JSON 模板"""
        logger.info("="*70)
        logger.info("🎨 PPTX 轉 JSON 模板工具")
        logger.info("="*70)
        
        # 1. 分析 PPTX
        analysis = self.analyze_pptx(pptx_path)
        
        # 保存分析結果（可選）
        analysis_path = pptx_path.replace('.pptx', '_analysis.json')
        with open(analysis_path, 'w', encoding='utf-8') as f:
            json.dump(analysis, f, ensure_ascii=False, indent=2)
        logger.info(f"\n💾 分析結果已保存：{analysis_path}")
        
        # 2. 使用 Gemini 生成模板
        template_json = self.generate_template_json(analysis, template_name)
        
        # 3. 保存模板
        if output_path is None:
            output_path = pptx_path.replace('.pptx', '_template.json')
        
        with open(output_path, 'w', encoding='utf-8') as f:
            # 格式化 JSON
            template_data = json.loads(template_json)
            json.dump(template_data, f, ensure_ascii=False, indent=2)
        
        logger.info(f"💾 模板已保存：{output_path}")
        
        # 4. 驗證模板
        self._validate_template(output_path)
        
        logger.info("\n" + "="*70)
        logger.info("✅ 完成！")
        logger.info("="*70)
        logger.info(f"\n生成的文件：")
        logger.info(f"  - 分析結果：{analysis_path}")
        logger.info(f"  - JSON 模板：{output_path}")
        
        return output_path
    
    def _validate_template(self, template_path: str):
        """驗證生成的模板"""
        logger.info(f"\n🔍 驗證模板...")
        
        try:
            with open(template_path, 'r', encoding='utf-8') as f:
                template = json.load(f)
            
            # 檢查必要字段
            assert 'template_info' in template, "缺少 template_info"
            assert 'slide_types' in template, "缺少 slide_types"
            
            info = template['template_info']
            assert 'name' in info, "template_info 缺少 name"
            assert 'slide_width' in info, "template_info 缺少 slide_width"
            assert 'slide_height' in info, "template_info 缺少 slide_height"
            
            slide_types = template['slide_types']
            logger.info(f"   ✓ 包含 {len(slide_types)} 種 slide 類型")
            
            for st in slide_types:
                assert 'type_id' in st, f"Slide type 缺少 type_id"
                assert 'pptx_layout' in st, f"Slide type {st.get('type_id')} 缺少 pptx_layout"
            
            logger.info(f"   ✓ 模板格式正確")
            
        except Exception as e:
            logger.error(f"   ❌ 模板驗證失敗：{e}")
            raise


def main():
    """測試主函數"""
    import os

    from dotenv import load_dotenv
    
    load_dotenv()
    
    api_key = os.getenv('GEMINI_API_KEY')
    if not api_key:
        print("❌ 請設置 GEMINI_API_KEY 環境變量")
        return
    
    # 創建轉換器
    creator = PPTXTemplateCreator(api_key=api_key)
    
    # 轉換 PPTX
    pptx_path = "pptx_template/test.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"❌ 文件不存在：{pptx_path}")
        return
    
    output_path = creator.create_template_from_pptx(
        pptx_path=pptx_path,
        output_path="templates/test_template.json",
        template_name="測試模板"
    )
    
    print(f"\n🎉 成功！模板已生成：{output_path}")


if __name__ == '__main__':
    main()
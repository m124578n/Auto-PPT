"""
Slide 類型定義（使用 Strategy Pattern + Registry Pattern）

設計理念：
1. 每個 slide 類型是一個獨立的策略類
2. 統一的介面：generate_html() 和 generate_pptx()
3. 使用 Registry 自動註冊所有 slide 類型
4. 新增 slide 類型只需繼承 SlideType 並實現方法即可
"""

import os
from typing import Any, Dict

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt

from AutoPPT.slide_types._base import SlideType
from AutoPPT.slide_types.slide_registry import SlideTypeRegistry
from AutoPPT.utils.logger import get_logger

# 获取日志器
logger = get_logger()

# ==================== 具體 Slide 類型 ====================

@SlideTypeRegistry.register('opening')
class OpeningSlide(SlideType):
    """開場頁"""
    
    @classmethod
    def get_json_example(cls) -> Dict[str, Any]:
        return {
            "slide_type": "opening",
            "title": "主標題",
            "subtitle": "副標題"
        }
    
    @classmethod
    def get_description(cls) -> str:
        return "開場頁（漸層背景）"
    
    def generate_html(self) -> str:
        title = self.data.get('title', '')
        subtitle = self.data.get('subtitle', '')
        
        return f"""
        <div class="slide slide-opening">
            <div class="slide-content">
                <h1 class="main-title">{title}</h1>
                {f'<p class="subtitle">{subtitle}</p>' if subtitle else ''}
            </div>
        </div>
        """
    
    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景漸層 - HTML: linear-gradient(135deg, #667eea 0%, #764ba2 100%)
        # 使用中間色 #6E72C6
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(110, 114, 198)
        bg.line.fill.background()
        
        # 主標題
        main_title = self.data.get('title', '')
        if main_title:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            
            # 處理換行符號，為每一行創建單獨的段落
            lines = main_title.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = title_frame.paragraphs[0]
                else:
                    p = title_frame.add_paragraph()
                
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據標題長度調整字體大小
                if len(main_title) > 15:
                    p.font.size = Pt(52)
                else:
                    p.font.size = Pt(58)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 副標題
        subtitle = self.data.get('subtitle', '')
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(1.2)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            
            # 處理換行符號，為每一行創建單獨的段落
            lines = subtitle.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = subtitle_frame.paragraphs[0]
                else:
                    p = subtitle_frame.add_paragraph()
                
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據副標題長度調整字體大小
                if len(subtitle) > 25:
                    p.font.size = Pt(26)
                else:
                    p.font.size = Pt(28)
                p.font.color.rgb = RGBColor(255, 255, 255)
        
        return slide


@SlideTypeRegistry.register('section_divider')
class SectionSlide(SlideType):
    """章節分隔頁"""
    
    @classmethod
    def get_json_example(cls) -> Dict[str, Any]:
        return {
            "slide_type": "section_divider",
            "section_title": "章節名稱"
        }
    
    @classmethod
    def get_description(cls) -> str:
        return "章節分隔頁（藍色背景）"
    
    def generate_html(self) -> str:
        section_title = self.data.get('section_title', '')
        
        return f"""
        <div class="slide slide-section">
            <div class="slide-content">
                <h2 class="section-title">{section_title}</h2>
                <div class="decoration-line"></div>
            </div>
        </div>
        """
    
    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景 - HTML: linear-gradient(135deg, #4682b4 0%, #2c5f8d 100%)
        # 使用中間色 #3970A1
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(57, 112, 161)
        bg.line.fill.background()
        
        # 章節標題
        section_title = self.data.get('section_title', '')
        if section_title:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.8)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            
            # 處理換行符號，為每一行創建單獨的段落
            lines = section_title.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = title_frame.paragraphs[0]
                else:
                    p = title_frame.add_paragraph()
                
                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據標題長度調整字體大小
                if len(section_title) > 20:
                    p.font.size = Pt(46)
                else:
                    p.font.size = Pt(50)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
        
        # 裝飾線（上）
        line_top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.6), Inches(2), Inches(0.04)
        )
        line_top.fill.solid()
        line_top.fill.fore_color.rgb = RGBColor(255, 255, 255)
        line_top.line.fill.background()
        
        # 裝飾線（下）
        if len(section_title) > 11:
            bottom = Inches(4.9)
        else:
            bottom = Inches(4.0)
        line_bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4), bottom, Inches(2), Inches(0.04)
        )
        line_bottom.fill.solid()
        line_bottom.fill.fore_color.rgb = RGBColor(255, 255, 255)
        line_bottom.line.fill.background()
        
        return slide


@SlideTypeRegistry.register('text_content')
class TextContentSlide(SlideType):
    """純文字內容頁"""

    @classmethod
    def get_json_example(cls) -> Dict[str, Any]:
        return {
            "slide_type": "text_content",
            "title": "頁面標題",
            "bullets": ["要點1", "要點2", "要點3"],
            "indent_levels": [0, 0, 1]
        }

    @classmethod
    def get_description(cls) -> str:
        return "純文字內容頁（項目符號列表，indent_levels: 0=主要點, 1=次要點），bullets要在五個以內，超過五個要分頁"

    def generate_html(self) -> str:
        title = self.data.get('title', '')
        bullets = self.data.get('bullets', [])
        indent_levels = self.data.get('indent_levels', [0] * len(bullets))

        bullets_html = ""
        for bullet, level in zip(bullets, indent_levels):
            indent_class = f"indent-{level}" if level > 0 else ""
            bullets_html += f'<li class="{indent_class}">{bullet}</li>\n'

        return f"""
        <div class="slide slide-content">
            <div class="slide-content">
                <h2 class="slide-title">{title}</h2>
                <ul class="bullet-list">
                    {bullets_html}
                </ul>
            </div>
        </div>
        """

    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 標題
        title = self.data.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.6), Inches(8.8), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            # 處理換行符號，為每一行創建單獨的段落
            lines = title.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = title_frame.paragraphs[0]
                else:
                    p = title_frame.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(38)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 62, 80)

            # 標題下劃線
            if len(title) < 17:
                bottom = 1.35
            else:
                bottom = 1.95
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.6),
                Inches(bottom),
                Inches(8.8),
                Inches(0.05),
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(70, 130, 180)
            line.line.fill.background()

        # 項目符號列表
        bullets = self.data.get('bullets', [])
        indent_levels = self.data.get('indent_levels', [0] * len(bullets))

        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(bottom + 0.45), Inches(8.0), Inches(5.4)
            )
            text_frame = content_box.text_frame
            text_frame.word_wrap = True

            # 根據項目數量調整字體
            num_items = len(bullets)
            avg_length = sum(len(b) for b in bullets) / num_items

            if num_items >= 5 and avg_length > 25:
                base_font_size, indent_font_size = 21, 19
                base_spacing, indent_spacing = 12, 10
                line_spacing = 1.25
            elif num_items >= 5:
                base_font_size, indent_font_size = 22, 20
                base_spacing, indent_spacing = 14, 12
                line_spacing = 1.3
            elif num_items >= 4:
                base_font_size, indent_font_size = 23, 21
                base_spacing, indent_spacing = 16, 14
                line_spacing = 1.35
            else:
                base_font_size, indent_font_size = 24, 22
                base_spacing, indent_spacing = 18, 16
                line_spacing = 1.4

            for i, (bullet, level) in enumerate(zip(bullets, indent_levels)):
                is_indent = level > 0

                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.line_spacing = line_spacing

                if is_indent:
                    # 次要項目
                    p.level = 1
                    p.space_after = Pt(indent_spacing)

                    # 使用 run 來分別設置箭頭和文字的大小
                    # 清空默認文本
                    if i == 0:
                        p.clear()

                    # 箭頭符號（固定大小）
                    run_bullet = p.add_run()
                    run_bullet.text = "▶ "
                    run_bullet.font.size = Pt(20)  # 箭頭固定大小
                    run_bullet.font.color.rgb = RGBColor(100, 100, 100)

                    # 文字內容（根據項目數量調整）
                    run_text = p.add_run()
                    run_text.text = bullet
                    run_text.font.size = Pt(indent_font_size)
                    run_text.font.color.rgb = RGBColor(85, 85, 85)
                else:
                    # 主要項目
                    p.level = 0
                    p.space_after = Pt(base_spacing)

                    # 使用 run 來分別設置箭頭和文字的大小
                    # 清空默認文本
                    if i == 0:
                        p.clear()

                    # 箭頭符號（固定大小）
                    run_bullet = p.add_run()
                    run_bullet.text = "▶ "
                    run_bullet.font.size = Pt(26)  # 箭頭固定大小
                    run_bullet.font.color.rgb = RGBColor(70, 130, 180)

                    # 文字內容（根據項目數量調整）
                    run_text = p.add_run()
                    run_text.text = bullet
                    run_text.font.size = Pt(base_font_size)
                    run_text.font.color.rgb = RGBColor(52, 73, 94)

        return slide


@SlideTypeRegistry.register('image_with_text')
class ImageTextSlide(SlideType):
    """圖文混合頁"""

    @classmethod
    def get_json_example(cls) -> Dict[str, Any]:
        return {
            "slide_type": "image_with_text",
            "title": "標題",
            "image_id": "img_01",
            "text": "說明文字",
            "layout": "horizontal"
        }

    @classmethod
    def get_description(cls) -> str:
        return "圖文混合頁（layout 可選 'horizontal' 左圖右文 或 'vertical' 上圖下文）"

    def generate_html(self) -> str:
        title = self.data.get('title', '')
        image_id = self.data.get('image_id', '')
        text = self.data.get('text', '')
        layout = self.data.get('layout', 'horizontal')

        img_src = self._get_image_path(image_id) or ""
        layout_class = "layout-vertical" if layout == "vertical" else "layout-horizontal"

        return f"""
        <div class="slide slide-content">
            <div class="slide-content">
                <h2 class="slide-title">{title}</h2>
                <div class="image-text-container {layout_class}">
                    <div class="image-box">
                        <img src="{img_src}" alt="">
                    </div>
                    <div class="text-box">
                        <p>{text.replace(chr(10), '<br>')}</p>
                    </div>
                </div>
            </div>
        </div>
        """

    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 標題
        title = self.data.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.75)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            # 處理換行符號，為每一行創建單獨的段落
            lines = title.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = title_frame.paragraphs[0]
                else:
                    p = title_frame.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據標題長度調整字體大小
                if len(title) > 20:
                    p.font.size = Pt(34)
                else:
                    p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50

        # 檢查布局類型
        layout = self.data.get('layout', 'horizontal')
        is_horizontal = layout == 'horizontal'

        # 提取圖片
        image_id = self.data.get('image_id', '')
        image_path = self._get_image_path(image_id)

        if image_path and os.path.exists(image_path):
            try:
                if is_horizontal:
                    # 左圖右文 - 保持寬高比
                    max_width = 4.4
                    max_height = 5.0  # 調整高度以適應整體布局
                    img_width, img_height = self._calculate_image_size(
                        image_path, max_width, max_height
                    )
                    # 圖片水平居中，垂直頂部對齊文字框
                    left = Inches(0.7 + (max_width - img_width) / 2)
                    top = Inches(2.5)  # 往下移，讓整體更平衡
                else:
                    # 上圖下文（vertical）- 保持寬高比
                    max_width = 8.0  # 與文字框寬度一致
                    max_height = 3.2
                    img_width, img_height = self._calculate_image_size(
                        image_path, max_width, max_height
                    )
                    # 圖片水平居中對齊
                    left = Inches(1.0 + (max_width - img_width) / 2)
                    top = Inches(2.5)  # 往下移，讓整體更平衡

                slide.shapes.add_picture(
                    image_path, left, top,
                    width=Inches(img_width),
                    height=Inches(img_height)
                )
                logger.info(
                    f'   ✓ 添加圖片：{os.path.basename(image_path)} ({img_width:.2f}" × {img_height:.2f}")'
                )
            except Exception as e:
                logger.info(f"   ⚠️ 圖片添加失敗：{e}")
        else:
            if image_id:
                logger.info(f"   ⚠️ 圖片不存在：{image_id}")

        # 提取文字
        text_content = self.data.get('text', '')
        if text_content:
            if is_horizontal:
                # 右側文字 - 擴大文字框以容納更多內容
                text_box = slide.shapes.add_textbox(
                    Inches(5.35), Inches(2.2), Inches(4.0), Inches(5.0)
                )
                # 根據文字長度調整字體大小
                if len(text_content) > 200:
                    font_size = 19
                    line_spacing = 1.5
                elif len(text_content) > 150:
                    font_size = 20
                    line_spacing = 1.55
                else:
                    font_size = 21
                    line_spacing = 1.6
            else:
                # 下方文字 - 擴大文字框高度
                text_box = slide.shapes.add_textbox(
                    Inches(1), Inches(5.5), Inches(8), Inches(1.7)
                )
                # 根據文字長度調整字體大小
                if len(text_content) > 150:
                    font_size = 18
                    line_spacing = 1.4
                else:
                    font_size = 19
                    line_spacing = 1.5

            text_frame = text_box.text_frame
            text_frame.text = text_content
            text_frame.word_wrap = True
            text_frame.paragraphs[0].font.size = Pt(font_size)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50
            text_frame.paragraphs[0].line_spacing = line_spacing

        return slide


@SlideTypeRegistry.register('full_image')
class FullImageSlide(SlideType):
    """大圖展示頁"""

    @classmethod
    def get_json_example(cls) -> Dict[str, Any]:
        return {
            "slide_type": "full_image",
            "title": "標題",
            "image_id": "img_02",
            "caption": "圖片說明"
        }

    @classmethod
    def get_description(cls) -> str:
        return "大圖展示頁（大幅圖片配簡短說明）"

    def generate_html(self) -> str:
        title = self.data.get('title', '')
        image_id = self.data.get('image_id', '')
        caption = self.data.get('caption', '')

        img_src = self._get_image_path(image_id) or ""

        return f"""
        <div class="slide slide-content">
            <div class="slide-content">
                <h2 class="slide-title">{title}</h2>
                <div class="full-image-container">
                    <img src="{img_src}" alt="">
                    {f'<p class="caption">{caption}</p>' if caption else ''}
                </div>
            </div>
        </div>
        """

    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 標題
        title = self.data.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(0.5), Inches(8), Inches(0.75)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            # 處理換行符號，為每一行創建單獨的段落
            lines = title.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = title_frame.paragraphs[0]
                else:
                    p = title_frame.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據標題長度調整字體大小
                if len(title) > 20:
                    p.font.size = Pt(34)
                else:
                    p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = RGBColor(44, 62, 80)  # #2c3e50

        # 提取圖片
        image_id = self.data.get('image_id', '')
        image_path = self._get_image_path(image_id)

        if image_path and os.path.exists(image_path):
            try:
                # 大圖展示 - HTML: max-width 90%, max-height 70%
                # 保持寬高比
                max_width = 7.5
                max_height = 4.2
                img_width, img_height = self._calculate_image_size(
                    image_path, max_width, max_height
                )
                # 居中對齊圖片
                left = Inches(1.25 + (max_width - img_width) / 2)
                top = Inches(1.9 + (max_height - img_height) / 2)

                slide.shapes.add_picture(
                    image_path, left, top,
                    width=Inches(img_width),
                    height=Inches(img_height)
                )
                logger.info(
                    f'   ✓ 添加大圖：{os.path.basename(image_path)} ({img_width:.2f}" × {img_height:.2f}")'
                )
            except Exception as e:
                logger.info(f"   ⚠️ 圖片添加失敗：{e}")
        else:
            if image_id:
                logger.info(f"   ⚠️ 圖片不存在：{image_id}")

        # 提取圖片說明
        caption = self.data.get('caption', '')
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(6.5), Inches(8.4), Inches(0.7)
            )
            caption_frame = caption_box.text_frame
            caption_frame.word_wrap = True

            # 處理換行符號，為每一行創建單獨的段落
            lines = caption.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = caption_frame.paragraphs[0]
                else:
                    p = caption_frame.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據說明文字長度調整字體大小
                if len(caption) > 60:
                    p.font.size = Pt(15)
                else:
                    p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(127, 140, 141)  # #7f8c8d

        return slide


@SlideTypeRegistry.register('closing')
class ClosingSlide(SlideType):
    """結尾頁"""

    @classmethod
    def get_json_example(cls) -> Dict[str, Any]:
        return {"slide_type": "closing", "closing_text": "謝謝觀看", "subtext": ""}

    @classmethod
    def get_description(cls) -> str:
        return "結尾頁（漸層背景）"

    def generate_html(self) -> str:
        closing_text = self.data.get('closing_text', '謝謝觀看')
        subtext = self.data.get('subtext', '')

        return f"""
        <div class="slide slide-closing">
            <div class="slide-content">
                <h1 class="closing-title">{closing_text}</h1>
                {f'<p class="closing-subtext">{subtext}</p>' if subtext else ''}
            </div>
        </div>
        """

    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 背景 - HTML: linear-gradient(135deg, #f093fb 0%, #f5576c 100%)
        # 使用中間色 #F375B4
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(243, 117, 180)
        bg.line.fill.background()

        # 結尾標題（注意：HTML中是 h1）
        closing_text = self.data.get('closing_text', '謝謝觀看')
        if closing_text:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(9), Inches(1.8)
            )
            title_frame = title_box.text_frame
            title_frame.word_wrap = True

            # 處理換行符號，為每一行創建單獨的段落
            lines = closing_text.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = title_frame.paragraphs[0]
                else:
                    p = title_frame.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據標題長度調整字體大小
                if len(closing_text) > 25:
                    p.font.size = Pt(46)
                else:
                    p.font.size = Pt(50)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

        # 提取副文字
        subtext = self.data.get('subtext', '')
        if subtext:
            subtext_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.5), Inches(9), Inches(1.0)
            )
            subtext_frame = subtext_box.text_frame
            subtext_frame.word_wrap = True

            # 處理換行符號，為每一行創建單獨的段落
            lines = subtext.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = subtext_frame.paragraphs[0]
                else:
                    p = subtext_frame.add_paragraph()

                p.text = line
                p.alignment = PP_ALIGN.CENTER
                # 根據副文字長度調整字體大小
                if len(subtext) > 20:
                    p.font.size = Pt(24)
                else:
                    p.font.size = Pt(26)
                p.font.color.rgb = RGBColor(255, 255, 255)

        return slide

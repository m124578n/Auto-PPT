"""
擴展範例：新增自定義 Slide 類型

這個範例展示如何輕鬆添加新的 Slide 類型，而不需要修改現有代碼。
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt

from AutoPPT.slide_types.slide_types import SlideType, SlideTypeRegistry


# ==================== 範例 1: 兩欄文字對比頁 ====================
@SlideTypeRegistry.register('two_column_text')
class TwoColumnTextSlide(SlideType):
    """兩欄文字對比頁
    
    適用場景：對比兩個概念、優缺點比較等
    """
    
    @classmethod
    def get_json_example(cls):
        """JSON 示例（會自動出現在 AI prompt 中）"""
        return {
            "slide_type": "two_column_text",
            "title": "標題",
            "left_title": "左欄標題",
            "left_content": "左欄內容",
            "right_title": "右欄標題",
            "right_content": "右欄內容"
        }
    
    @classmethod
    def get_description(cls):
        """類型說明（會自動出現在 AI prompt 中）"""
        return "兩欄文字對比頁（適合對比兩個概念、優缺點比較等）"
    
    def generate_html(self) -> str:
        title = self.data.get('title', '')
        left_title = self.data.get('left_title', '左欄')
        left_content = self.data.get('left_content', '')
        right_title = self.data.get('right_title', '右欄')
        right_content = self.data.get('right_content', '')
        
        return f"""
        <div class="slide slide-content">
            <div class="slide-content">
                <h2 class="slide-title">{title}</h2>
                <div style="display: grid; grid-template-columns: 1fr 1fr; gap: 40px; flex: 1;">
                    <div style="background: #f8f9fa; padding: 30px; border-radius: 10px;">
                        <h3 style="color: #4682b4; margin-bottom: 20px; font-size: 28px;">{left_title}</h3>
                        <p style="font-size: 22px; color: #2c3e50; line-height: 1.6;">{left_content}</p>
                    </div>
                    <div style="background: #f8f9fa; padding: 30px; border-radius: 10px;">
                        <h3 style="color: #4682b4; margin-bottom: 20px; font-size: 28px;">{right_title}</h3>
                        <p style="font-size: 22px; color: #2c3e50; line-height: 1.6;">{right_content}</p>
                    </div>
                </div>
            </div>
        </div>
        """
    
    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 標題
        title = self.data.get('title', '')
        if title:
            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.6), Inches(8.8), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(38)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
            
            # 標題下劃線
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(8.8), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(70, 130, 180)
            line.line.fill.background()
        
        # 左欄背景
        left_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.8), Inches(4.2), Inches(5)
        )
        left_bg.fill.solid()
        left_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        left_bg.line.color.rgb = RGBColor(200, 200, 200)
        
        # 右欄背景
        right_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.2), Inches(5)
        )
        right_bg.fill.solid()
        right_bg.fill.fore_color.rgb = RGBColor(248, 249, 250)
        right_bg.line.color.rgb = RGBColor(200, 200, 200)
        
        # 左欄標題
        left_title = self.data.get('left_title', '左欄')
        left_title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2), Inches(3.8), Inches(0.6)
        )
        left_title_frame = left_title_box.text_frame
        left_title_frame.text = left_title
        left_title_frame.paragraphs[0].font.size = Pt(26)
        left_title_frame.paragraphs[0].font.bold = True
        left_title_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
        
        # 左欄內容
        left_content = self.data.get('left_content', '')
        left_content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.8), Inches(3.8), Inches(3.8)
        )
        left_content_frame = left_content_box.text_frame
        left_content_frame.text = left_content
        left_content_frame.word_wrap = True
        left_content_frame.paragraphs[0].font.size = Pt(20)
        left_content_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        left_content_frame.paragraphs[0].line_spacing = 1.5
        
        # 右欄標題
        right_title = self.data.get('right_title', '右欄')
        right_title_box = slide.shapes.add_textbox(
            Inches(5.4), Inches(2), Inches(3.8), Inches(0.6)
        )
        right_title_frame = right_title_box.text_frame
        right_title_frame.text = right_title
        right_title_frame.paragraphs[0].font.size = Pt(26)
        right_title_frame.paragraphs[0].font.bold = True
        right_title_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
        
        # 右欄內容
        right_content = self.data.get('right_content', '')
        right_content_box = slide.shapes.add_textbox(
            Inches(5.4), Inches(2.8), Inches(3.8), Inches(3.8)
        )
        right_content_frame = right_content_box.text_frame
        right_content_frame.text = right_content
        right_content_frame.word_wrap = True
        right_content_frame.paragraphs[0].font.size = Pt(20)
        right_content_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        right_content_frame.paragraphs[0].line_spacing = 1.5
        
        return slide


# ==================== 範例 2: 引用卡片頁 ====================
@SlideTypeRegistry.register('quote_card')
class QuoteCardSlide(SlideType):
    """引用卡片頁
    
    適用場景：展示名言、客戶評價等
    """
    
    @classmethod
    def get_json_example(cls):
        """JSON 示例（會自動出現在 AI prompt 中）"""
        return {
            "slide_type": "quote_card",
            "quote": "引用內容",
            "author": "作者",
            "author_title": "作者職稱"
        }
    
    @classmethod
    def get_description(cls):
        """類型說明（會自動出現在 AI prompt 中）"""
        return "引用卡片頁（適合展示名言、客戶評價、重要引述等）"
    
    def generate_html(self) -> str:
        quote = self.data.get('quote', '')
        author = self.data.get('author', '')
        author_title = self.data.get('author_title', '')
        
        return f"""
        <div class="slide slide-quote">
            <div class="slide-content" style="display: flex; align-items: center; justify-content: center;">
                <div style="max-width: 800px; text-align: center;">
                    <div style="font-size: 48px; color: #4682b4; margin-bottom: 20px;">"</div>
                    <p style="font-size: 32px; color: #2c3e50; line-height: 1.6; font-style: italic; margin-bottom: 40px;">
                        {quote}
                    </p>
                    <div style="font-size: 48px; color: #4682b4; margin-bottom: 20px;">"</div>
                    <div style="border-top: 3px solid #4682b4; padding-top: 30px;">
                        <p style="font-size: 24px; color: #34495e; font-weight: bold; margin-bottom: 10px;">
                            {author}
                        </p>
                        {f'<p style="font-size: 20px; color: #7f8c8d;">{author_title}</p>' if author_title else ''}
                    </div>
                </div>
            </div>
        </div>
        """
    
    def generate_pptx(self, prs: Presentation) -> Slide:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景裝飾
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(7), Inches(4.5)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(248, 249, 250)
        bg_shape.line.color.rgb = RGBColor(70, 130, 180)
        bg_shape.line.width = Pt(3)
        
        # 開引號
        open_quote_box = slide.shapes.add_textbox(
            Inches(2), Inches(1.8), Inches(1), Inches(0.8)
        )
        open_quote_frame = open_quote_box.text_frame
        open_quote_frame.text = '"'
        open_quote_frame.paragraphs[0].font.size = Pt(72)
        open_quote_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
        
        # 引用內容
        quote = self.data.get('quote', '')
        quote_box = slide.shapes.add_textbox(
            Inches(2.2), Inches(2.5), Inches(5.6), Inches(2)
        )
        quote_frame = quote_box.text_frame
        quote_frame.text = quote
        quote_frame.word_wrap = True
        quote_frame.paragraphs[0].font.size = Pt(28)
        quote_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)
        quote_frame.paragraphs[0].font.italic = True
        quote_frame.paragraphs[0].line_spacing = 1.5
        
        # 閉引號
        close_quote_box = slide.shapes.add_textbox(
            Inches(7), Inches(4.2), Inches(1), Inches(0.8)
        )
        close_quote_frame = close_quote_box.text_frame
        close_quote_frame.text = '"'
        close_quote_frame.paragraphs[0].font.size = Pt(72)
        close_quote_frame.paragraphs[0].font.color.rgb = RGBColor(70, 130, 180)
        
        # 分隔線
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(2.2), Inches(4.8), Inches(5.6), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(70, 130, 180)
        line.line.fill.background()
        
        # 作者
        author = self.data.get('author', '')
        author_box = slide.shapes.add_textbox(
            Inches(2.2), Inches(5.1), Inches(5.6), Inches(0.5)
        )
        author_frame = author_box.text_frame
        author_frame.text = author
        author_frame.paragraphs[0].font.size = Pt(22)
        author_frame.paragraphs[0].font.bold = True
        author_frame.paragraphs[0].font.color.rgb = RGBColor(52, 73, 94)
        
        # 作者職稱
        author_title = self.data.get('author_title', '')
        if author_title:
            author_title_box = slide.shapes.add_textbox(
                Inches(2.2), Inches(5.6), Inches(5.6), Inches(0.4)
            )
            author_title_frame = author_title_box.text_frame
            author_title_frame.text = author_title
            author_title_frame.paragraphs[0].font.size = Pt(18)
            author_title_frame.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)
        
        return slide


# ==================== 測試新的 Slide 類型 ====================
def test_new_slide_types():
    """測試新增的 Slide 類型"""
    import json

    from slide_generator import HTMLGenerator, PPTXGenerator

    print("=" * 60)
    print("測試新增的 Slide 類型")
    print("=" * 60)

    # 測試數據
    test_data = {
        'title': '新 Slide 類型示範',
        'topic': 'new_slide_types_demo',
        'slides': [
            {
                'slide_type': 'opening',
                'title': '新 Slide 類型示範',
                'subtitle': '展示如何輕鬆擴展功能'
            },
            {
                'slide_type': 'two_column_text',
                'title': '傳統方法 vs 新架構',
                'left_title': '傳統方法 ❌',
                'left_content': '需要修改多個文件，HTML 和 PPTX 生成邏輯分散，難以維護，添加新功能成本高。',
                'right_title': '新架構 ✅',
                'right_content': '只需添加一個類，HTML 和 PPTX 邏輯統一，易於維護，擴展簡單快速。'
            },
            {
                'slide_type': 'quote_card',
                'quote': '好的架構設計能讓你的代碼更易讀、更易維護、更易擴展。',
                'author': '資深工程師',
                'author_title': 'Software Architect'
            },
            {
                'slide_type': 'text_content',
                'title': '如何添加新 Slide 類型',
                'bullets': [
                    '繼承 SlideType 基類',
                    '使用 @SlideTypeRegistry.register() 裝飾器',
                    '實現 generate_html() 方法',
                    '實現 generate_pptx() 方法',
                    '完成！立即可用'
                ],
                'indent_levels': [0, 0, 1, 1, 0]
            },
            {
                'slide_type': 'closing',
                'closing_text': '開始創建你的 Slide 類型吧！',
                'subtext': '發揮你的創意'
            }
        ]
    }

    # 顯示已註冊的類型
    print(f"已註冊的 Slide 類型：{SlideTypeRegistry.all_types()}")
    print(f"✓ 成功註冊了 {len(SlideTypeRegistry.all_types())} 種類型")

    # 生成 HTML
    print("生成 HTML...")
    html_gen = HTMLGenerator()
    html_content = html_gen.generate_from_data(test_data)

    html_file = 'example_new_slide_types.html'
    with open(html_file, 'w', encoding='utf-8') as f:
        f.write(html_content)
    print(f"✓ HTML 已保存：{html_file}")

    # 生成 PPTX
    print("生成 PPTX...")
    pptx_gen = PPTXGenerator()
    prs = pptx_gen.generate_from_data(test_data)

    pptx_file = 'example_new_slide_types.pptx'
    pptx_gen.save(pptx_file)
    print(f"✓ PPTX 已保存：{pptx_file}")

    # 保存 JSON
    json_file = 'example_new_slide_types_data.json'
    with open(json_file, 'w', encoding='utf-8') as f:
        json.dump(test_data, f, ensure_ascii=False, indent=2)
    print(f"✓ JSON 已保存：{json_file}")

    print("=" * 60)
    print("✅ 測試完成！")
    print("=" * 60)
    print("💡 重點：")
    print("  - 只需添加新類並註冊，不需修改其他代碼")
    print("  - HTML 和 PPTX 生成邏輯集中在一個類中")
    print("  - 自動整合到現有工作流程")
    print()


if __name__ == "__main__":
    test_new_slide_types()

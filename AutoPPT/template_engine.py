"""
模板引擎 - 動態加載和管理 PPTX 模板

核心功能：
1. 從 JSON 文件加載模板定義
2. 動態生成 SlideType 類
3. 生成 AI Prompt
4. 管理 Slide 佈局和樣式
"""

import json
import os
from dataclasses import dataclass
from typing import Any, Dict, List, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt

from AutoPPT.utils.logger import get_logger

logger = get_logger()


@dataclass
class Position:
    """位置信息"""
    left: float
    top: float
    width: float
    height: float
    
    @classmethod
    def from_dict(cls, data: Dict) -> 'Position':
        """從字典創建（支持 width/height 或 max_width/max_height）"""
        return cls(
            left=data.get('left', 0),
            top=data.get('top', 0),
            width=data.get('width', data.get('max_width', 0)),
            height=data.get('height', data.get('max_height', 0))
        )
    
    def to_inches(self):
        """轉換為 Inches"""
        return (Inches(self.left), Inches(self.top), 
                Inches(self.width), Inches(self.height))


@dataclass
class ElementStyle:
    """元素樣式"""
    data: Dict
    
    def get(self, key: str, default=None):
        """獲取樣式屬性"""
        return self.data.get(key, default)
    
    def get_color_rgb(self, key: str, default="#000000") -> RGBColor:
        """獲取 RGB 顏色"""
        color_hex = self.data.get(key, default)
        color_hex = color_hex.lstrip('#')
        r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)


@dataclass
class SlideElement:
    """Slide 元素定義"""
    type: str  # textbox, image, shape
    name: str
    position: Optional[Position] = None
    position_horizontal: Optional[Position] = None
    position_vertical: Optional[Position] = None
    style: Optional[ElementStyle] = None
    shape_type: Optional[str] = None
    
    @classmethod
    def from_dict(cls, data: Dict) -> 'SlideElement':
        """從字典創建"""
        position = None
        position_horizontal = None
        position_vertical = None
        
        if 'position' in data:
            position = Position.from_dict(data['position'])
        if 'position_horizontal' in data:
            position_horizontal = Position.from_dict(data['position_horizontal'])
        if 'position_vertical' in data:
            position_vertical = Position.from_dict(data['position_vertical'])
        
        style = ElementStyle(data.get('style', {})) if 'style' in data else None
        
        return cls(
            type=data['type'],
            name=data['name'],
            position=position,
            position_horizontal=position_horizontal,
            position_vertical=position_vertical,
            style=style,
            shape_type=data.get('shape_type')
        )


@dataclass
class SlideTypeDefinition:
    """Slide 類型定義"""
    type_id: str
    name: str
    description: str
    llm_instruction: str
    json_schema: Dict
    layout_index: int
    background: Dict
    elements: List[SlideElement]
    
    @classmethod
    def from_dict(cls, data: Dict) -> 'SlideTypeDefinition':
        """從字典創建"""
        layout = data['pptx_layout']
        elements = [SlideElement.from_dict(e) for e in layout.get('elements', [])]
        
        return cls(
            type_id=data['type_id'],
            name=data['name'],
            description=data['description'],
            llm_instruction=data['llm_instruction'],
            json_schema=data['json_schema'],
            layout_index=layout.get('layout_index', 6),
            background=layout.get('background', {}),
            elements=elements
        )


class PPTXTemplate:
    """PPTX 模板管理器"""

    def __init__(self, json_path: str = None, pptx_path: str = None):
        """
        初始化模板

        Args:
            json_path: JSON 配置文件路徑
            pptx_path: PPTX 模板文件路徑（可選，用於保留原始設計）
        """
        self.json_path = json_path or self._get_default_template_path()
        self.pptx_path = pptx_path
        self.template_data = {}
        self.slide_types: Dict[str, SlideTypeDefinition] = {}
        self.pptx_template = None  # 存儲 PPTX 模板對象

        # 加載模板
        self._load_template()

        # 如果提供了 PPTX 模板，加載它
        if self.pptx_path:
            self._load_pptx_template()

    def _get_default_template_path(self) -> str:
        """獲取默認模板路徑"""
        return os.path.join(
            os.path.dirname(os.path.dirname(__file__)),
            'templates', 'default_template.json'
        )

    def _load_template(self):
        """加載模板配置文件"""
        try:
            logger.info(f"📄 加載模板配置：{self.json_path}")

            with open(self.json_path, "r", encoding="utf-8") as f:
                self.template_data = json.load(f)

            # 解析 Slide 類型
            for slide_data in self.template_data.get('slide_types', []):
                slide_def = SlideTypeDefinition.from_dict(slide_data)
                self.slide_types[slide_def.type_id] = slide_def

            template_info = self.template_data.get('template_info', {})
            logger.info(f"   ✓ 模板：{template_info.get('name', 'Unknown')}")
            logger.info(f"   ✓ 版本：{template_info.get('version', 'Unknown')}")
            logger.info(f"   ✓ Slide 類型數量：{len(self.slide_types)}")

        except FileNotFoundError:
            logger.error(f"❌ 模板配置文件不存在：{self.json_path}")
            raise
        except json.JSONDecodeError as e:
            logger.error(f"❌ 模板 JSON 格式錯誤：{e}")
            raise

    def _load_pptx_template(self):
        """加載 PPTX 模板文件"""
        try:
            logger.info(f"📄 加載 PPTX 模板：{self.pptx_path}")

            from pptx import Presentation

            self.pptx_template = Presentation(self.pptx_path)

            logger.info(f"   ✓ PPTX 模板加載成功")
            logger.info(f"   ✓ 可用布局數量：{len(self.pptx_template.slide_layouts)}")

            # 列出所有可用的布局
            for i, layout in enumerate(self.pptx_template.slide_layouts):
                logger.info(f"      - Layout {i}: {layout.name}")

        except FileNotFoundError:
            logger.error(f"❌ PPTX 模板文件不存在：{self.pptx_path}")
            raise
        except Exception as e:
            logger.error(f"❌ 加載 PPTX 模板失敗：{e}")
            raise

    def get_slide_type_definition(self, type_id: str) -> Optional[SlideTypeDefinition]:
        """獲取 Slide 類型定義"""
        return self.slide_types.get(type_id)

    def get_all_slide_type_ids(self) -> List[str]:
        """獲取所有 Slide 類型 ID"""
        return list(self.slide_types.keys())

    def generate_ai_prompt(self, image_metadata: Dict = None, user_prompt: str = "") -> str:
        """
        生成 AI Prompt
        
        Args:
            image_metadata: 圖片元數據
            user_prompt: 用戶提示詞
            
        Returns:
            完整的 AI Prompt
        """
        # 圖片列表信息
        image_list_info = "無圖片資源（純文字簡報）"
        if image_metadata:
            image_list_info = "\n".join([
                f"- {img_id}: {data['filename']}"
                for img_id, data in image_metadata.items()
            ])

        # 生成 JSON Schema 示例
        json_examples = []
        for type_id, slide_def in self.slide_types.items():
            example = slide_def.json_schema.copy()
            json_examples.append(example)

        slides_examples_str = ",\n    ".join([
            json.dumps(example, ensure_ascii=False, indent=2).replace("\n", "\n    ")
            for example in json_examples
        ])

        # 生成類型說明
        descriptions = []
        for type_id, slide_def in self.slide_types.items():
            descriptions.append(
                f"- {type_id}: {slide_def.name} - {slide_def.llm_instruction}"
            )
        descriptions_str = "\n".join(descriptions)

        # 構建完整 Prompt
        prompt = f"""請分析以下內容，生成一個結構化的演示文稿。

**使用者輸入**
{user_prompt}

**文字內容**：
請讀取我上傳的檔案，當作其內容。

**可用圖片**：
{image_list_info}

**輸出 JSON 格式**：
{{
  "title": "簡報標題",
  "topic": "簡報主題",
  "slides": [
    {slides_examples_str}
  ]
}}

**可用的 Slide 類型說明**：
{descriptions_str}

**要求**：
1. 自動分析內容，識別2-4個主題
2. 每個主題有章節分隔頁
3. 合理安排圖片（如有）
4. 總共10-15張幻燈片
5. 嚴格按照上述 JSON 格式輸出
6. 避免使用markdown格式
"""
        return prompt

    def get_presentation_config(self) -> Dict:
        """獲取 Presentation 配置"""
        info = self.template_data.get('template_info', {})
        return {
            'slide_width': info.get('slide_width', 10.0),
            'slide_height': info.get('slide_height', 7.5),
        }

    def create_slide(
        self,
        prs: Presentation,
        slide_data: Dict,
        image_metadata: Dict = None
    ) -> Slide:
        """
        根據數據創建 Slide
        
        Args:
            prs: Presentation 對象
            slide_data: Slide 數據
            image_metadata: 圖片元數據
            
        Returns:
            創建的 Slide
        """
        slide_type_id = slide_data.get('slide_type', 'text_content')
        slide_def = self.get_slide_type_definition(slide_type_id)

        if not slide_def:
            logger.warning(f"⚠️  未知的 Slide 類型：{slide_type_id}")
            slide_def = self.get_slide_type_definition('text_content')

        # 創建 slide
        layout_index = slide_def.layout_index

        # 使用 prs 的布局（如果有 PPTX 模板，prs 就是從模板創建的）
        if layout_index < len(prs.slide_layouts):
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            logger.info(
                f"   ✓ 使用布局 {layout_index}: {prs.slide_layouts[layout_index].name}"
            )
        else:
            logger.warning(f"⚠️  布局索引 {layout_index} 超出範圍，使用布局 0")
            slide = prs.slides.add_slide(prs.slide_layouts[0])

        # 檢查是否使用 PPTX 模板
        # 只有在沒有 PPTX 模板時才添加背景和元素
        # 因為 PPTX 模板已經包含了設計
        if not self.pptx_path:
            # 沒有 PPTX 模板，手動添加背景和元素
            self._add_background(slide, prs, slide_def.background)

            # 添加元素
            for element in slide_def.elements:
                self._add_element(slide, element, slide_data, image_metadata)
        else:
            # 使用 PPTX 模板時，填充佔位符
            self._fill_placeholders(slide, slide_data, image_metadata)

        return slide

    def _fill_placeholders(
        self, slide: Slide, slide_data: Dict, image_metadata: Dict = None
    ):
        """填充 PPTX 模板的佔位符"""
        try:
            filled_count = 0

            # 準備要填充的內容
            title = slide_data.get("title") or slide_data.get("section_title") or ""
            subtitle = slide_data.get("subtitle") or slide_data.get("subtext") or ""
            image_id = slide_data.get("image_id")

            # 收集所有內容字段（按順序）
            content_fields = []
            for key in ["content", "text", "closing_text", "content_1", "content_2"]:
                value = slide_data.get(key)
                if value:
                    content_fields.append((key, value))

            bullets = slide_data.get("bullets", [])

            logger.info(
                f"      數據：標題={bool(title)}, 副標題={bool(subtitle)}, 內容={len(content_fields)}, 圖片={bool(image_id)}, Bullets={len(bullets)}"
            )

            # 按佔位符類型填充
            content_index = 0
            for shape in slide.shapes:
                try:
                    if not hasattr(shape, "placeholder_format"):
                        continue

                    placeholder_type = shape.placeholder_format.type

                    # 標題佔位符 (TITLE=1, CENTER_TITLE=3)
                    if placeholder_type in [1, 3]:
                        if title:
                            shape.text = title
                            logger.info(f"      ✓ 填充標題：{title[:40]}...")
                            filled_count += 1

                    # 副標題佔位符 (SUBTITLE=4)
                    elif placeholder_type == 4:
                        if subtitle:
                            shape.text = subtitle
                            logger.info(f"      ✓ 填充副標題：{subtitle[:40]}...")
                            filled_count += 1

                    # 內容佔位符 (BODY=2, OBJECT=7)
                    elif placeholder_type in [2, 7]:
                        # 優先填充 bullets
                        if bullets:
                            text_frame = shape.text_frame
                            text_frame.clear()

                            for bullet, level in zip(
                                bullets,
                                slide_data.get("indent_levels", [0] * len(bullets)),
                            ):
                                p = text_frame.add_paragraph()
                                p.text = bullet
                                p.level = level

                            logger.info(f"      ✓ 填充 Bullets：{len(bullets)} 個")
                            filled_count += 1
                            bullets = []  # 清空，避免重複填充

                        # 其次填充文本內容
                        elif content_index < len(content_fields):
                            key, content = content_fields[content_index]
                            content = self._clean_markdown(str(content))

                            if "\n" in content:
                                # 多行文本
                                text_frame = shape.text_frame
                                text_frame.clear()
                                for i, line in enumerate(content.split("\n")):
                                    line = line.strip()
                                    if line:
                                        if i == 0:
                                            p = text_frame.paragraphs[0]
                                        else:
                                            p = text_frame.add_paragraph()
                                        p.text = line
                            else:
                                # 單行文本
                                shape.text = content

                            logger.info(f"      ✓ 填充內容 {key}：{content[:40]}...")
                            filled_count += 1
                            content_index += 1

                        # 最後嘗試填充副標題（如果還沒有被 SUBTITLE 佔位符填充）
                        elif subtitle and placeholder_type == 2:  # 只在 BODY 類型
                            shape.text = subtitle
                            logger.info(
                                f"      ✓ 填充副標題（BODY）：{subtitle[:40]}..."
                            )
                            filled_count += 1
                            subtitle = ""  # 清空，避免重複填充

                    # 圖片佔位符 (PICTURE=18)
                    elif placeholder_type == 18:
                        if image_id and image_metadata and image_id in image_metadata:
                            image_path = image_metadata[image_id].get("path")
                            if image_path and os.path.exists(image_path):
                                try:
                                    shape.insert_picture(image_path)
                                    logger.info(f"      ✓ 填充圖片：{image_id}")
                                    filled_count += 1
                                except Exception as img_error:
                                    logger.warning(
                                        f"      ⚠️  圖片插入失敗：{img_error}"
                                    )

                except Exception as shape_error:
                    logger.debug(f"      處理 shape 時出錯：{shape_error}")
                    continue

            if filled_count == 0:
                logger.warning(f"      ⚠️  沒有填充任何內容！請檢查佔位符設置。")
            else:
                logger.info(f"      ✓ 總共填充了 {filled_count} 個元素")

        except Exception as e:
            logger.error(f"❌ 填充佔位符失敗：{e}")
            import traceback

            logger.error(traceback.format_exc())

    def _add_background(self, slide: Slide, prs: Presentation, bg_config: Dict):
        """添加背景"""
        if not bg_config:
            return

        bg_type = bg_config.get('type', 'solid')

        if bg_type == 'solid':
            color_hex = bg_config.get('color', '#FFFFFF')
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            bg.fill.solid()

            color_hex = color_hex.lstrip('#')
            r, g, b = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
            bg.fill.fore_color.rgb = RGBColor(r, g, b)
            bg.line.fill.background()

        elif bg_type == 'gradient':
            # 使用中間色作為近似
            start_hex = bg_config.get('color_start', '#FFFFFF').lstrip('#')
            end_hex = bg_config.get('color_end', '#FFFFFF').lstrip('#')

            start_rgb = tuple(int(start_hex[i:i+2], 16) for i in (0, 2, 4))
            end_rgb = tuple(int(end_hex[i:i+2], 16) for i in (0, 2, 4))

            # 計算中間色
            mid_rgb = tuple((s + e) // 2 for s, e in zip(start_rgb, end_rgb))

            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(*mid_rgb)
            bg.line.fill.background()

    def _add_element(
        self,
        slide: Slide,
        element: SlideElement,
        slide_data: Dict,
        image_metadata: Dict = None
    ):
        """添加元素"""
        if element.type == 'textbox':
            self._add_textbox(slide, element, slide_data)
        elif element.type == 'image':
            self._add_image(slide, element, slide_data, image_metadata)
        elif element.type == 'shape':
            self._add_shape(slide, element)

    def _add_textbox(self, slide: Slide, element: SlideElement, slide_data: Dict):
        """添加文本框"""
        # 獲取位置（支持 position、position_horizontal、position_vertical）
        position = element.position

        # 如果有 horizontal/vertical 位置，根據 layout 選擇
        if not position and (element.position_horizontal or element.position_vertical):
            layout = slide_data.get('layout', 'horizontal')
            if layout == 'horizontal' and element.position_horizontal:
                position = element.position_horizontal
            elif layout == 'vertical' and element.position_vertical:
                position = element.position_vertical
            elif element.position_horizontal:  # 默認使用 horizontal
                position = element.position_horizontal

        if not position:
            return

        # 檢查是否是 bullets 內容（特殊處理）
        is_bullets = element.name == 'content' and 'bullets' in slide_data

        # 如果不是 bullets，檢查文本值
        if not is_bullets:
            text_value = slide_data.get(element.name, '')
            if not text_value:
                return

        left, top, width, height = position.to_inches()
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # 處理特殊情況：bullets
        if is_bullets:
            self._add_bullet_content(text_frame, slide_data, element.style)
        else:
            # 普通文本
            text_value = slide_data.get(element.name, '')
            # 清理 markdown 格式
            text_value = self._clean_markdown(str(text_value))

            # 處理換行符
            if '\n' in text_value:
                self._add_multiline_text(text_frame, text_value, element, slide_data)
            else:
                p = text_frame.paragraphs[0]
                p.text = text_value
                self._apply_text_style(p, element, slide_data)

    def _clean_markdown(self, text: str) -> str:
        """清理 markdown 格式"""
        import re

        # 移除粗體標記 **text**
        text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
        # 移除斜體標記 *text*
        text = re.sub(r'\*([^*]+)\*', r'\1', text)
        # 移除其他常見 markdown 標記
        text = re.sub(r'__([^_]+)__', r'\1', text)
        text = re.sub(r'_([^_]+)_', r'\1', text)
        return text

    def _add_multiline_text(self, text_frame, text_value: str, element: SlideElement, slide_data: Dict):
        """添加多行文本（處理換行符）"""
        lines = text_value.split('\n')

        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue

            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = line
            self._apply_text_style(p, element, slide_data)

            # 設置行間距
            if element.style and element.style.get('line_spacing'):
                p.line_spacing = element.style.get('line_spacing')

    def _apply_text_style(self, paragraph, element: SlideElement, slide_data: Dict):
        """應用文本樣式"""
        if not element.style:
            return

        # 對齊方式
        if element.style.get('alignment') == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif element.style.get('alignment') == 'left':
            paragraph.alignment = PP_ALIGN.LEFT
        elif element.style.get('alignment') == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT

        # 字體大小（支持 horizontal/vertical 特定大小）
        layout = slide_data.get('layout', 'horizontal')
        font_size_key = f'font_size_{layout}' if layout in ['horizontal', 'vertical'] else 'font_size'
        font_size = element.style.get(font_size_key, element.style.get('font_size'))

        if font_size and paragraph.runs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
        elif font_size:
            paragraph.font.size = Pt(font_size)

        # 粗體
        if element.style.get('font_bold'):
            if paragraph.runs:
                for run in paragraph.runs:
                    run.font.bold = True
            else:
                paragraph.font.bold = True

        # 顏色
        font_color = element.style.get('font_color')
        if font_color:
            color_rgb = element.style.get_color_rgb('font_color')
            if paragraph.runs:
                for run in paragraph.runs:
                    run.font.color.rgb = color_rgb
            else:
                paragraph.font.color.rgb = color_rgb

    def _add_bullet_content(self, text_frame, slide_data: Dict, style: ElementStyle = None):
        """添加項目符號內容"""
        bullets = slide_data.get('bullets', [])
        indent_levels = slide_data.get('indent_levels', [0] * len(bullets))

        if not bullets:
            return

        # 獲取樣式（如果 style 為 None，使用默認值）
        bullet_symbol_base = style.get('bullet_symbol_base', '▸') if style else '▸'
        bullet_symbol_indent = style.get('bullet_symbol_indent', '▸') if style else '▸'
        font_size_base = style.get('font_size_base', 24) if style else 24
        font_size_indent = style.get('font_size_indent', 22) if style else 22
        bullet_size_base = style.get('bullet_size_base', 26) if style else 26
        bullet_size_indent = style.get('bullet_size_indent', 20) if style else 20

        for i, (bullet, level) in enumerate(zip(bullets, indent_levels)):
            is_indent = level > 0

            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.line_spacing = 1.3

            if i == 0:
                p.clear()

            # 添加箭頭符號
            run_bullet = p.add_run()
            if is_indent:
                run_bullet.text = f"{bullet_symbol_indent} "
                run_bullet.font.size = Pt(bullet_size_indent)
                if style:
                    run_bullet.font.color.rgb = style.get_color_rgb('bullet_color_indent', '#646464')
                else:
                    run_bullet.font.color.rgb = RGBColor(100, 100, 100)
            else:
                run_bullet.text = f"{bullet_symbol_base} "
                run_bullet.font.size = Pt(bullet_size_base)
                if style:
                    run_bullet.font.color.rgb = style.get_color_rgb('bullet_color_base', '#4682B4')
                else:
                    run_bullet.font.color.rgb = RGBColor(70, 130, 180)

            # 添加文字
            run_text = p.add_run()
            run_text.text = bullet
            if is_indent:
                run_text.font.size = Pt(font_size_indent)
                if style:
                    run_text.font.color.rgb = style.get_color_rgb('font_color_indent', '#555555')
                else:
                    run_text.font.color.rgb = RGBColor(85, 85, 85)
                p.level = 1
            else:
                run_text.font.size = Pt(font_size_base)
                if style:
                    run_text.font.color.rgb = style.get_color_rgb('font_color_base', '#34495E')
                else:
                    run_text.font.color.rgb = RGBColor(52, 73, 94)
                p.level = 0

    def _add_image(
        self,
        slide: Slide,
        element: SlideElement,
        slide_data: Dict,
        image_metadata: Dict = None
    ):
        """添加圖片"""
        image_id = slide_data.get('image_id', '')
        if not image_id or not image_metadata:
            return

        image_path = None
        if image_id in image_metadata:
            image_path = image_metadata[image_id].get('path')

        if not image_path or not os.path.exists(image_path):
            logger.warning(f"⚠️  圖片不存在：{image_id}")
            return

        # 根據 layout 選擇位置
        layout = slide_data.get('layout', 'horizontal')
        if layout == 'vertical' and element.position_vertical:
            position = element.position_vertical
        elif element.position_horizontal:
            position = element.position_horizontal
        else:
            position = element.position

        if not position:
            return

        # 計算圖片尺寸（保持寬高比）
        from PIL import Image
        try:
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height

                max_width = position.width
                max_height = position.height

                # 計算實際尺寸
                if aspect_ratio >= max_width / max_height:
                    # 寬度優先
                    actual_width = max_width
                    actual_height = max_width / aspect_ratio
                else:
                    # 高度優先
                    actual_height = max_height
                    actual_width = max_height * aspect_ratio

                # 居中
                left = position.left + (max_width - actual_width) / 2
                top = position.top + (max_height - actual_height) / 2

                slide.shapes.add_picture(
                    image_path,
                    Inches(left),
                    Inches(top),
                    width=Inches(actual_width),
                    height=Inches(actual_height)
                )
                logger.info(f"   ✓ 添加圖片：{os.path.basename(image_path)}")

        except Exception as e:
            logger.warning(f"⚠️  添加圖片失敗：{e}")

    def _add_shape(self, slide: Slide, element: SlideElement):
        """添加形狀"""
        if not element.position or not element.shape_type:
            return

        left, top, width, height = element.position.to_inches()

        if element.shape_type == 'rectangle':
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()

            if element.style:
                fill_color = element.style.get('fill_color')
                if fill_color:
                    shape.fill.fore_color.rgb = element.style.get_color_rgb('fill_color')

            shape.line.fill.background()

    def __repr__(self) -> str:
        """字符串表示"""
        info = self.template_data.get('template_info', {})
        return f"<PPTXTemplate: {info.get('name', 'Unknown')} v{info.get('version', '?')}>"
